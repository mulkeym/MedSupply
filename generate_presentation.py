"""Generate MedSupply POC overview presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MED_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
TEAL = RGBColor(0x0D, 0x94, 0x88)


def add_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=DARK_TEXT,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_frame(slide, left, top, width, height, items, font_size=14,
                     color=DARK_TEXT, spacing=Pt(6), bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        if bold_first and i == 0:
            p.font.bold = True
        # Handle bold prefix with "**prefix** rest" syntax
        if "**" in item:
            parts = item.split("**")
            for j, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                run.text = part
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.name = "Calibri"
                run.font.bold = (j % 2 == 1)
        else:
            p.text = item
    return txBox


def add_section_header(slide, text):
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.15), NAVY)
    add_text_box(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
                 text, font_size=32, color=WHITE, bold=True)
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(1.0), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()


def add_footer(slide, slide_num, total):
    add_text_box(slide, Inches(0.6), Inches(7.0), Inches(6), Inches(0.4),
                 "Medical Resupply Predictor POC", font_size=10, color=MED_GRAY)
    add_text_box(slide, Inches(11), Inches(7.0), Inches(2), Inches(0.4),
                 f"{slide_num} / {total}", font_size=10, color=MED_GRAY,
                 alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 12

# ════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2),
             "Medical Resupply Predictor",
             font_size=48, color=WHITE, bold=True)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.8), Inches(2.85), Inches(3), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.8),
             "Proof of Concept Overview",
             font_size=28, color=ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.0),
             "Phlebotomy-Support Supply Risk Prediction\nAcross a Notional USINDOPACOM Medical Logistics Network",
             font_size=18, color=MED_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
             "Simulated, Unclassified Data",
             font_size=14, color=MED_GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 2: Problem Statement
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Problem Statement")

# Left column - The Challenge
card = add_shape(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), LIGHT_GRAY, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.65), Inches(5.2), Inches(0.5),
             "The Challenge", font_size=20, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(0.9), Inches(2.2), Inches(5.2), Inches(4.2), [
    "Military medical logistics networks must maintain readiness across geographically dispersed nodes",
    "Operational events (MASCAL, FPCON, hub disruptions) create sudden demand surges and supply-chain disruptions",
    "Planners lack forward-looking visibility into which nodes will experience stockouts and when",
    "Manual planning cannot efficiently evaluate compound scenarios across a multi-echelon supply network",
    "Without predictive tools, resupply decisions are reactive rather than proactive",
], font_size=14, color=DARK_TEXT)

# Right column - Key Questions
card = add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), LIGHT_BLUE, ACCENT_BLUE)
add_text_box(slide, Inches(7.1), Inches(1.65), Inches(5.2), Inches(0.5),
             "Key Questions to Answer", font_size=20, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(4.2), [
    "Which nodes fall below reorder point?",
    "Which nodes stock out during the forecast horizon?",
    "Which item causes the first mission risk?",
    "How do MASCAL or FPCON events affect demand?",
    "How does hub or route disruption affect resupply latency?",
    "Which secondary or tertiary route should be activated?",
    "What minimum supply levels should be targeted?",
], font_size=14, color=DARK_TEXT)

add_footer(slide, 2, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 3: POC Scope
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "POC Scope & Approach")

# Supply items card
card = add_shape(slide, Inches(0.6), Inches(1.5), Inches(3.7), Inches(5.2), LIGHT_GRAY)
add_text_box(slide, Inches(0.9), Inches(1.65), Inches(3.2), Inches(0.5),
             "8 Phlebotomy Supplies", font_size=18, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(0.9), Inches(2.2), Inches(3.2), Inches(4.0), [
    "Nitrile gloves (high criticality)",
    "Blood collection tubes (high)",
    "Butterfly needle sets (high)",
    "Alcohol prep pads (medium)",
    "Gauze pads (medium)",
    "Tourniquets (medium)",
    "Specimen bags (medium)",
    "Barcode labels (low)",
], font_size=13, color=DARK_TEXT)

# Network card
card = add_shape(slide, Inches(4.6), Inches(1.5), Inches(3.7), Inches(5.2), LIGHT_GRAY)
add_text_box(slide, Inches(4.9), Inches(1.65), Inches(3.2), Inches(0.5),
             "14-Node Network", font_size=18, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(4.9), Inches(2.2), Inches(3.2), Inches(4.0), [
    "1 Strategic Supplier",
    "1 DLA Prime Vendor",
    "1 Theater Hub",
    "3 Regional Hubs",
    "3 MTFs (Medical Treatment Facilities)",
    "3 Forward Medical Units",
    "2 Battalion Aid Stations",
    "19 supply routes (primary / secondary / tertiary)",
], font_size=13, color=DARK_TEXT)

# Scenarios card
card = add_shape(slide, Inches(8.6), Inches(1.5), Inches(4.1), Inches(5.2), LIGHT_GRAY)
add_text_box(slide, Inches(8.9), Inches(1.65), Inches(3.5), Inches(0.5),
             "5 Event Types", font_size=18, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(8.9), Inches(2.2), Inches(3.5), Inches(4.0), [
    "**MASCAL** - mass casualty demand surge",
    "**FPCON** - force protection restrictions",
    "**Hub Degraded** - route delays & reliability loss",
    "**Hub Offline** - primary routes blocked",
    "**Inventory Loss** - stock reduction at node",
    "",
    "Each event has low / medium / high severity",
    "Events queue for compound scenarios",
    "1-180 day forecast horizons",
], font_size=13, color=DARK_TEXT)

add_footer(slide, 3, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 4: Architecture
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "System Architecture")

# Browser box
box = add_shape(slide, Inches(0.6), Inches(1.6), Inches(4.0), Inches(5.0), LIGHT_BLUE, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.75), Inches(3.4), Inches(0.5),
             "Browser Frontend", font_size=20, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(0.9), Inches(2.35), Inches(3.4), Inches(3.8), [
    "Google Maps regional supply-chain map",
    "Node selection & status coloring",
    "Operations tab (map, events, forecasts)",
    "Data tab (node directory, inventory editing)",
    "Event queue & scenario controls",
    "Impact summary & risk tables",
    "Global stock-posture scenarios",
], font_size=13, color=DARK_TEXT)

# Arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(4.9), Inches(3.8), Inches(0.8), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT_BLUE
arrow.line.fill.background()

# Backend box
box = add_shape(slide, Inches(6.0), Inches(1.6), Inches(4.0), Inches(5.0), LIGHT_GRAY, NAVY)
add_text_box(slide, Inches(6.3), Inches(1.75), Inches(3.4), Inches(0.5),
             "FastAPI Backend", font_size=20, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(6.3), Inches(2.35), Inches(3.4), Inches(3.8), [
    "REST API (17 endpoints)",
    "Per-node burn-rate engine",
    "Compound event processing",
    "Route impact & reachability analysis",
    "Path latency calculations",
    "Inventory forecasting",
    "Recommendation generation",
], font_size=13, color=DARK_TEXT)

# Arrow to DB
arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  Inches(10.3), Inches(3.8), Inches(0.8), Inches(0.5))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = ACCENT_BLUE
arrow2.line.fill.background()

# DB box
box = add_shape(slide, Inches(11.3), Inches(2.2), Inches(1.6), Inches(3.6), NAVY)
add_text_box(slide, Inches(11.35), Inches(2.35), Inches(1.5), Inches(0.5),
             "SQLite", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_frame(slide, Inches(11.35), Inches(2.9), Inches(1.5), Inches(2.8), [
    "Items",
    "Nodes",
    "Routes",
    "Op States",
    "Demand Profiles",
    "Inventory",
], font_size=11, color=LIGHT_BLUE, spacing=Pt(4))

add_footer(slide, 4, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 5: Network Topology
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Network Topology")

# Describe the hub-and-spoke model
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
             "Directed graph modeling a USINDOPACOM-style hub-and-spoke medical supply chain",
             font_size=16, color=DARK_TEXT)

# Draw simplified network diagram using shapes
echelons = [
    ("Strategic Supplier", 6.2, 2.2, NAVY),
    ("DLA Prime Vendor", 6.2, 3.0, DARK_BLUE),
    ("Pacific Theater Hub", 6.2, 3.8, ACCENT_BLUE),
]
hubs = [
    ("Regional Hub\nNorth", 2.8, 4.7, TEAL),
    ("Regional Hub\nCentral", 6.2, 4.7, TEAL),
    ("Regional Hub\nSouth", 9.6, 4.7, TEAL),
]
leaves = [
    ("MTF Alpha", 1.2, 5.6, GREEN),
    ("Fwd Red", 0.2, 6.4, AMBER),
    ("BAS Silver", 2.2, 6.4, AMBER),
    ("MTF Bravo", 5.2, 5.6, GREEN),
    ("Fwd Blue", 5.2, 6.4, AMBER),
    ("MTF Charlie", 8.6, 5.6, GREEN),
    ("BAS Gold", 8.0, 6.4, AMBER),
    ("Fwd Green", 9.8, 6.4, AMBER),
]

# Draw boxes
def draw_node(slide, label, x, y, color, w=1.8, h=0.55):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

for label, x, y, color in echelons:
    draw_node(slide, label, x, y, color)
for label, x, y, color in hubs:
    draw_node(slide, label, x, y, color, w=1.8, h=0.6)
for label, x, y, color in leaves:
    draw_node(slide, label, x, y, color, w=1.5, h=0.5)

# Route legend
card = add_shape(slide, Inches(11.0), Inches(2.0), Inches(2.0), Inches(2.8), LIGHT_GRAY)
add_text_box(slide, Inches(11.1), Inches(2.05), Inches(1.8), Inches(0.4),
             "Route Types", font_size=13, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_frame(slide, Inches(11.1), Inches(2.45), Inches(1.8), Inches(2.0), [
    "Primary (1-14 days)",
    "Secondary (2-7 days)",
    "Tertiary (7 days)",
    "",
    "19 total routes",
    "Latency: 1-14 days",
    "Reliability: 0.68-0.93",
], font_size=10, color=DARK_TEXT, spacing=Pt(3))

# Node type legend
card = add_shape(slide, Inches(11.0), Inches(5.0), Inches(2.0), Inches(1.8), LIGHT_GRAY)
add_text_box(slide, Inches(11.1), Inches(5.05), Inches(1.8), Inches(0.4),
             "Node Types", font_size=13, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_frame(slide, Inches(11.1), Inches(5.4), Inches(1.8), Inches(1.3), [
    "Supplier / DLA",
    "Theater Hub",
    "Regional Hubs (3)",
    "MTFs (3)",
    "Forward Units (3)",
    "Battalion Aid (2)",
], font_size=10, color=DARK_TEXT, spacing=Pt(2))

add_footer(slide, 5, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 6: Data Model
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Data Model")

# Table: Items
card = add_shape(slide, Inches(0.6), Inches(1.5), Inches(3.7), Inches(2.4), LIGHT_GRAY)
add_text_box(slide, Inches(0.9), Inches(1.55), Inches(3.2), Inches(0.4),
             "items", font_size=16, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(0.9), Inches(1.95), Inches(3.2), Inches(1.8), [
    "id (PK), name, unit",
    "usage_per_draw, usage_rate",
    "demand_basis (phlebotomy_event, specimen, patient_encounter, personnel_day)",
    "criticality (high, medium, low)",
], font_size=11, color=DARK_TEXT, spacing=Pt(3))

# Table: Nodes
card = add_shape(slide, Inches(4.6), Inches(1.5), Inches(3.7), Inches(2.4), LIGHT_GRAY)
add_text_box(slide, Inches(4.9), Inches(1.55), Inches(3.2), Inches(0.4),
             "nodes", font_size=16, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(4.9), Inches(1.95), Inches(3.2), Inches(1.8), [
    "id (PK), name, type",
    "latitude, longitude, x, y",
    "population, optempo",
    "stock_days",
], font_size=11, color=DARK_TEXT, spacing=Pt(3))

# Table: Routes
card = add_shape(slide, Inches(8.6), Inches(1.5), Inches(4.1), Inches(2.4), LIGHT_GRAY)
add_text_box(slide, Inches(8.9), Inches(1.55), Inches(3.5), Inches(0.4),
             "routes", font_size=16, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(8.9), Inches(1.95), Inches(3.5), Inches(1.8), [
    "id (PK), source_node_id, destination_node_id",
    "priority (primary, secondary, tertiary)",
    "days (delivery latency)",
    "reliability (0.0 - 1.0)",
], font_size=11, color=DARK_TEXT, spacing=Pt(3))

# Table: Operational States
card = add_shape(slide, Inches(0.6), Inches(4.2), Inches(3.7), Inches(2.6), LIGHT_GRAY)
add_text_box(slide, Inches(0.9), Inches(4.25), Inches(3.2), Inches(0.4),
             "operational_states", font_size=16, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(0.9), Inches(4.65), Inches(3.2), Inches(2.0), [
    "id (PK), label, description",
    "demand_multiplier (0.35x - 3.0x)",
    "route_latency_multiplier (1.0x - 1.35x)",
    "7 states from Standby to Combat Ops",
], font_size=11, color=DARK_TEXT, spacing=Pt(3))

# Table: Node Demand Profiles
card = add_shape(slide, Inches(4.6), Inches(4.2), Inches(3.7), Inches(2.6), LIGHT_GRAY)
add_text_box(slide, Inches(4.9), Inches(4.25), Inches(3.2), Inches(0.4),
             "node_demand_profiles", font_size=16, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(4.9), Inches(4.65), Inches(3.2), Inches(2.0), [
    "node_id (PK)",
    "active_supported_population",
    "daily_encounter_rate",
    "phlebotomy_probability",
    "specimens_per_phlebotomy",
    "operational_state, waste_factor",
], font_size=11, color=DARK_TEXT, spacing=Pt(3))

# Table: Inventory Balances
card = add_shape(slide, Inches(8.6), Inches(4.2), Inches(4.1), Inches(2.6), LIGHT_GRAY)
add_text_box(slide, Inches(8.9), Inches(4.25), Inches(3.5), Inches(0.4),
             "inventory_balances", font_size=16, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(8.9), Inches(4.65), Inches(3.5), Inches(2.0), [
    "node_id + item_id (composite PK)",
    "quantity_on_hand",
    "Initialized from burn-rate profile * stock_days * item_skew",
    "Hub balances sized from aggregate downstream demand",
], font_size=11, color=DARK_TEXT, spacing=Pt(3))

add_footer(slide, 6, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 7: Demand Model
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Demand Model")

# Formula card
card = add_shape(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(2.5), NAVY)
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(7.0), Inches(0.4),
             "Core Daily Demand Formula", font_size=20, color=WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(7.0), Inches(1.7))
tf = txBox.text_frame
tf.word_wrap = True
lines = [
    "daily_demand  =  active_supported_population",
    "                         x  workload_driver",
    "                         x  item_usage_rate",
    "                         x  operational_state_multiplier",
    "                         x  event_demand_multiplier",
    "                         x  waste_factor",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Consolas"
    p.space_after = Pt(1)

# Workload driver card
card = add_shape(slide, Inches(8.5), Inches(1.5), Inches(4.2), Inches(2.5), LIGHT_BLUE, ACCENT_BLUE)
add_text_box(slide, Inches(8.8), Inches(1.6), Inches(3.6), Inches(0.4),
             "Workload Drivers", font_size=18, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(8.8), Inches(2.1), Inches(3.6), Inches(1.7), [
    "**Phlebotomy event:** pop x encounter rate x phlebotomy prob",
    "**Specimen:** above x specimens per phlebotomy",
    "**Patient encounter:** pop x encounter rate",
    "**Personnel day:** pop x usage rate",
], font_size=12, color=DARK_TEXT, spacing=Pt(4))

# Operational states table
card = add_shape(slide, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.6), LIGHT_GRAY)
add_text_box(slide, Inches(0.9), Inches(4.35), Inches(5.2), Inches(0.4),
             "7 Operational States", font_size=18, color=NAVY, bold=True)

states_data = [
    ("Reduced Manning / Standby", "0.35x", "1.00x"),
    ("Garrison / Steady State", "1.00x", "1.00x"),
    ("Training / Exercise", "1.50x", "1.00x"),
    ("Deployment Prep / Mobilization", "1.75x", "1.05x"),
    ("Forward Deployed", "1.75x", "1.15x"),
    ("Active Operations", "2.00x", "1.20x"),
    ("Combat Operations", "3.00x", "1.35x"),
]

y_start = 4.85
for i, (state, demand, latency) in enumerate(states_data):
    color = DARK_TEXT
    add_text_box(slide, Inches(0.9), Inches(y_start + i * 0.24), Inches(2.8), Inches(0.24),
                 state, font_size=10, color=color)
    add_text_box(slide, Inches(3.7), Inches(y_start + i * 0.24), Inches(1.0), Inches(0.24),
                 demand, font_size=10, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.8), Inches(y_start + i * 0.24), Inches(1.0), Inches(0.24),
                 latency, font_size=10, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# Labels
add_text_box(slide, Inches(3.7), Inches(4.7), Inches(1.0), Inches(0.2),
             "Demand", font_size=9, color=MED_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.8), Inches(4.7), Inches(1.0), Inches(0.2),
             "Latency", font_size=9, color=MED_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# Default parameters
card = add_shape(slide, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.6), LIGHT_GRAY)
add_text_box(slide, Inches(7.1), Inches(4.35), Inches(5.2), Inches(0.4),
             "Default Profile Parameters", font_size=18, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(7.1), Inches(4.85), Inches(5.2), Inches(1.8), [
    "**Daily encounter rate:** 0.018 (1.8% of population per day)",
    "**Phlebotomy probability:** 0.38 (38% of encounters need a draw)",
    "**Specimens per phlebotomy:** 1.2",
    "**Waste factor:** 1.08 (8% waste / damage allowance)",
    "**Initial stock posture:** 180 days of baseline supply",
    "",
    "Demand is modeled per-node, not globally. Each node has its own demand profile.",
], font_size=12, color=DARK_TEXT, spacing=Pt(3))

add_footer(slide, 7, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 8: Event Model
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Event Model")

add_text_box(slide, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5),
             "Events are temporary scenario overlays. Multiple events can be queued for compound scenarios.",
             font_size=15, color=DARK_TEXT)

# MASCAL card
card = add_shape(slide, Inches(0.6), Inches(2.0), Inches(2.3), Inches(4.5), LIGHT_GRAY, RED)
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(2.1), Inches(0.4),
             "MASCAL", font_size=18, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_frame(slide, Inches(0.7), Inches(2.55), Inches(2.1), Inches(3.5), [
    "Demand Effect:",
    "  Days 1-3: peak surge",
    "  Day 4+: sustained surge",
    "",
    "Low: 3x / 1.5x",
    "Med: 5x / 2x",
    "High: 8x / 3x",
    "",
    "No direct supply effect",
    "Cascades downstream",
], font_size=11, color=DARK_TEXT, spacing=Pt(2))

# FPCON card
card = add_shape(slide, Inches(3.15), Inches(2.0), Inches(2.3), Inches(4.5), LIGHT_GRAY, AMBER)
add_text_box(slide, Inches(3.25), Inches(2.1), Inches(2.1), Inches(0.4),
             "FPCON", font_size=18, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_frame(slide, Inches(3.25), Inches(2.55), Inches(2.1), Inches(3.5), [
    "Demand Effect:",
    "  Low: 1.10x",
    "  Med: 1.25x",
    "  High: 1.45x",
    "",
    "Supply Effect:",
    "  Route latency increase",
    "  Low: 1.25x",
    "  Med: 1.50x",
    "  High: 2.00x",
], font_size=11, color=DARK_TEXT, spacing=Pt(2))

# Hub Degraded card
card = add_shape(slide, Inches(5.7), Inches(2.0), Inches(2.3), Inches(4.5), LIGHT_GRAY, TEAL)
add_text_box(slide, Inches(5.8), Inches(2.1), Inches(2.1), Inches(0.4),
             "Hub Degraded", font_size=18, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_frame(slide, Inches(5.8), Inches(2.55), Inches(2.1), Inches(3.5), [
    "No demand increase",
    "",
    "Supply Effect:",
    "  Outbound latency doubled",
    "  Reliability reduced",
    "",
    "  Low: -12% reliability",
    "  Med: -25% reliability",
    "  High: -40% reliability",
], font_size=11, color=DARK_TEXT, spacing=Pt(2))

# Hub Offline card
card = add_shape(slide, Inches(8.25), Inches(2.0), Inches(2.3), Inches(4.5), LIGHT_GRAY, NAVY)
add_text_box(slide, Inches(8.35), Inches(2.1), Inches(2.1), Inches(0.4),
             "Hub Offline", font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_frame(slide, Inches(8.35), Inches(2.55), Inches(2.1), Inches(3.5), [
    "No demand increase",
    "",
    "Supply Effect:",
    "  All primary outbound",
    "  routes BLOCKED",
    "",
    "  Downstream nodes",
    "  fall to organic",
    "  stock levels",
    "  (14-120 days by type)",
], font_size=11, color=DARK_TEXT, spacing=Pt(2))

# Inventory Loss card
card = add_shape(slide, Inches(10.8), Inches(2.0), Inches(2.0), Inches(4.5), LIGHT_GRAY, DARK_BLUE)
add_text_box(slide, Inches(10.85), Inches(2.1), Inches(1.9), Inches(0.4),
             "Inv. Loss", font_size=18, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_frame(slide, Inches(10.85), Inches(2.55), Inches(1.9), Inches(3.5), [
    "No demand increase",
    "",
    "Inventory reduced:",
    "  Low: 25% lost",
    "  Med: 50% lost",
    "  High: 75% lost",
    "",
    "Immediate effect",
    "at selected node",
], font_size=11, color=DARK_TEXT, spacing=Pt(2))

add_footer(slide, 8, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 9: Inventory Policy
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Inventory Policy & Thresholds")

# Thresholds
card = add_shape(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.2), LIGHT_GRAY)
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5.2), Inches(0.4),
             "Three Inventory Thresholds", font_size=20, color=NAVY, bold=True)

thresh = [
    ("Target Stock", "30 days of baseline demand", GREEN),
    ("Reorder Point", "21 days of baseline demand", AMBER),
    ("Critical Level", "7 days of baseline demand", RED),
]
for i, (label, desc, color) in enumerate(thresh):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.9), Inches(2.2 + i * 0.5), Inches(0.3), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_box(slide, Inches(1.35), Inches(2.2 + i * 0.5), Inches(1.8), Inches(0.35),
                 label, font_size=14, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(3.2), Inches(2.2 + i * 0.5), Inches(2.8), Inches(0.35),
                 desc, font_size=14, color=MED_GRAY)

# Net drawdown model
card = add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.2), NAVY)
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.4),
             "Net Drawdown Model", font_size=20, color=WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.15), Inches(5.2), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "net_drawdown = event_demand - baseline_replenishment"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT_BLUE
p.font.name = "Consolas"
p2 = tf.add_paragraph()
p2.text = ""
p2.font.size = Pt(6)
p3 = tf.add_paragraph()
p3.text = "In steady state (no events), net drawdown = 0 for reachable nodes."
p3.font.size = Pt(13)
p3.font.color.rgb = LIGHT_BLUE
p4 = tf.add_paragraph()
p4.text = "Events create risk by increasing demand, reducing inventory, or disrupting routes."
p4.font.size = Pt(13)
p4.font.color.rgb = LIGHT_BLUE

# Forecast status
card = add_shape(slide, Inches(0.6), Inches(4.0), Inches(5.8), Inches(2.6), LIGHT_BLUE)
add_text_box(slide, Inches(0.9), Inches(4.1), Inches(5.2), Inches(0.4),
             "Forecast Status Classification", font_size=20, color=NAVY, bold=True)

statuses = [
    ("Healthy", "> reorder point", GREEN),
    ("Watch", "< reorder point", AMBER),
    ("At Risk", "Forecasted below critical level", RGBColor(0xEA, 0x58, 0x0C)),
    ("Stockout", "Forecasted to reach zero", RED),
    ("Offline", "Hub offline event active", NAVY),
]
for i, (label, desc, color) in enumerate(statuses):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(0.9), Inches(4.65 + i * 0.38), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text_box(slide, Inches(1.25), Inches(4.6 + i * 0.38), Inches(1.3), Inches(0.3),
                 label, font_size=13, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(2.6), Inches(4.6 + i * 0.38), Inches(3.2), Inches(0.3),
                 desc, font_size=13, color=MED_GRAY)

# Organic fallback
card = add_shape(slide, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.6), LIGHT_GRAY)
add_text_box(slide, Inches(7.1), Inches(4.1), Inches(5.2), Inches(0.4),
             "Organic Stock Fallback (Unreachable Nodes)", font_size=18, color=NAVY, bold=True)

fallback = [
    ("Theater hub", "120 days"),
    ("Regional hub", "60 days"),
    ("MTF", "45 days"),
    ("Forward unit", "21 days"),
    ("Battalion aid", "14 days"),
]
for i, (ntype, days) in enumerate(fallback):
    add_text_box(slide, Inches(7.1), Inches(4.65 + i * 0.35), Inches(2.5), Inches(0.3),
                 ntype, font_size=13, color=DARK_TEXT)
    add_text_box(slide, Inches(9.6), Inches(4.65 + i * 0.35), Inches(1.5), Inches(0.3),
                 days, font_size=13, color=ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(7.1), Inches(6.3), Inches(5.2), Inches(0.4),
             "When all routes to a node are blocked, it falls back to organic stock posture.",
             font_size=11, color=MED_GRAY)

add_footer(slide, 9, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 10: Prediction Pipeline
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "How Predictions Are Created")

add_text_box(slide, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5),
             "The forecast engine runs a deterministic, day-by-day simulation for each node across the forecast horizon.",
             font_size=15, color=DARK_TEXT)

# Step boxes
steps = [
    ("1", "Load State", "Load items, nodes, routes,\nprofiles, states, and\ncurrent inventory from DB"),
    ("2", "Process Events", "Parse event queue into\nseverity parameters.\nApply inventory loss events."),
    ("3", "Route Impacts", "Calculate route status:\nnormal, delayed, blocked,\nor alternate active"),
    ("4", "Reachability", "Determine which nodes\ncan receive supply via\nnon-blocked routes"),
    ("5", "Path Latency", "Compute shortest-path\nlatency; compare baseline\nvs. current for gap days"),
    ("6", "Daily Simulation", "For each node and item,\nsimulate daily net burn\nacross the forecast horizon"),
]

for i, (num, title, desc) in enumerate(steps):
    x = 0.6 + i * 2.1
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(x + 0.7), Inches(2.05), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Card
    card = add_shape(slide, Inches(x), Inches(2.7), Inches(1.9), Inches(2.2), LIGHT_GRAY)
    add_text_box(slide, Inches(x + 0.1), Inches(2.75), Inches(1.7), Inches(0.35),
                 title, font_size=13, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.1), Inches(3.15), Inches(1.7), Inches(1.5),
                 desc, font_size=10, color=DARK_TEXT)

# Detail box for daily simulation
card = add_shape(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.6), NAVY)
add_text_box(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.35),
             "Step 6 Detail: Daily Inventory Simulation", font_size=16, color=WHITE, bold=True)
txBox = slide.shapes.add_textbox(Inches(0.9), Inches(5.7), Inches(5.5), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
code_lines = [
    "for day in range(1, horizon + 1):",
    "    event_demand = supported_daily_demand(events, day)",
    "    replenishment = baseline if reachable & day > latency_gap",
    "    net_burn = max(0, event_demand - replenishment)",
    "    balance -= net_burn",
    "    if balance <= 0: record stockout_day",
]
for i, line in enumerate(code_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(10)
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Consolas"
    p.space_after = Pt(0)

add_bullet_frame(slide, Inches(6.8), Inches(5.7), Inches(5.5), Inches(1.0), [
    "Supported demand includes hub's own demand + all primary downstream nodes",
    "Replenishment resumes only after path latency gap expires",
    "Stockout day, critical day, and remaining days of supply tracked per item",
    "Node status derived from worst-case item projection",
], font_size=10, color=LIGHT_BLUE, spacing=Pt(2))

add_footer(slide, 10, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 11: Prediction Outputs
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Prediction Outputs & Recommendations")

# Node-level outputs
card = add_shape(slide, Inches(0.6), Inches(1.5), Inches(3.7), Inches(3.2), LIGHT_GRAY)
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(3.2), Inches(0.4),
             "Per-Node Results", font_size=18, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(0.9), Inches(2.05), Inches(3.2), Inches(2.4), [
    "Node risk status (healthy/watch/risk/stockout/offline)",
    "First stockout day",
    "First critical day",
    "Riskiest item",
    "Effective stock days",
    "Supply reachability status",
    "Path latency (baseline vs. current)",
], font_size=12, color=DARK_TEXT, spacing=Pt(3))

# Item-level outputs
card = add_shape(slide, Inches(4.6), Inches(1.5), Inches(3.7), Inches(3.2), LIGHT_GRAY)
add_text_box(slide, Inches(4.9), Inches(1.6), Inches(3.2), Inches(0.4),
             "Per-Item Results", font_size=18, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(4.9), Inches(2.05), Inches(3.2), Inches(2.4), [
    "Current balance",
    "Reorder point (21-day demand)",
    "Critical level (7-day demand)",
    "Stockout day",
    "Critical day",
    "Remaining days of supply",
    "Daily burn rate / net burn rate",
    "Baseline replenishment rate",
], font_size=12, color=DARK_TEXT, spacing=Pt(3))

# Route outputs
card = add_shape(slide, Inches(8.6), Inches(1.5), Inches(4.1), Inches(3.2), LIGHT_GRAY)
add_text_box(slide, Inches(8.9), Inches(1.6), Inches(3.5), Inches(0.4),
             "Route & Network Results", font_size=18, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(8.9), Inches(2.05), Inches(3.5), Inches(2.4), [
    "Route status (normal/delayed/blocked/alternate)",
    "Current latency vs. baseline",
    "Reliability after degradation",
    "Network-wide metrics:",
    "  - Nodes at risk count",
    "  - Earliest stockout (any node)",
    "  - Routes impacted count",
], font_size=12, color=DARK_TEXT, spacing=Pt(3))

# Recommendations
card = add_shape(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.7), LIGHT_BLUE, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.4),
             "Automated Recommendations", font_size=20, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(0.9), Inches(5.55), Inches(11.5), Inches(1.0), [
    "**Alternate routes:** \"Activate secondary route from [hub]\" when non-primary routes exist for at-risk nodes",
    "**Emergency packages:** \"Push a 72-hour emergency phlebotomy package\" for forward units and battalion aid stations without viable routes",
    "**Unreachable nodes:** \"No viable upstream route; push emergency package or restore a supply line\" when all routes are blocked",
    "**Offline nodes:** \"Restore node operations or shift supported demand to alternate hubs\" for hub-offline events",
], font_size=12, color=DARK_TEXT, spacing=Pt(3))

add_footer(slide, 11, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 12: Summary & Future Work
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "Summary & Future Direction")

# What the POC demonstrates
card = add_shape(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(3.0), LIGHT_GRAY)
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5.2), Inches(0.4),
             "What This POC Demonstrates", font_size=20, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(0.9), Inches(2.1), Inches(5.2), Inches(2.2), [
    "Rule-based demand forecasting for transparency and auditability",
    "Multi-echelon supply chain modeling with hub-and-spoke topology",
    "Compound scenario analysis (multiple events, multiple nodes)",
    "Day-by-day inventory simulation with route disruption effects",
    "Actionable recommendations for supply-chain planners",
    "Interactive map-based visualization with Google Maps",
    "Editable demand profiles, node details, and inventory",
], font_size=13, color=DARK_TEXT, spacing=Pt(3))

# Future work
card = add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.0), LIGHT_BLUE, ACCENT_BLUE)
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.4),
             "High-Value Next Steps", font_size=20, color=NAVY, bold=True)
add_bullet_frame(slide, Inches(7.1), Inches(2.1), Inches(5.2), Inches(2.2), [
    "CSV import/export for items, nodes, routes, and inventory",
    "Store and compare forecast runs and scenario results",
    "Route capacity and partial fulfillment logic",
    "Item substitution rules",
    "Shelf-life and expiration-aware inventory",
    "User roles and audit history",
    "Replace synthetic data with authoritative logistics feeds",
], font_size=13, color=DARK_TEXT, spacing=Pt(3))

# Production direction
card = add_shape(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.7), NAVY)
add_text_box(slide, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.4),
             "Production Architecture Direction", font_size=20, color=WHITE, bold=True)
add_bullet_frame(slide, Inches(0.9), Inches(5.4), Inches(5.2), Inches(1.0), [
    "PostgreSQL (+ PostGIS for geospatial overlays)",
    "Scheduled ingestion jobs for external supply and encounter data",
    "Audit tables for scenario runs and user changes",
], font_size=13, color=LIGHT_BLUE, spacing=Pt(3))
add_bullet_frame(slide, Inches(6.8), Inches(5.4), Inches(5.2), Inches(1.0), [
    "Datasets needed: inventory, item catalogs, population data, encounter volume, requisition history, fulfillment lead times, route capacity",
    "Quality metrics: demand forecast MAE, stockout prediction precision/recall, days-of-supply error",
], font_size=13, color=LIGHT_BLUE, spacing=Pt(3))

add_footer(slide, 12, TOTAL_SLIDES)

# ── Save ──
output_path = "/Users/michaelmulkey/Documents/Repositories/MedSupply/MedSupply_POC_Overview.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
